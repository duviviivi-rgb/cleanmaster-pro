#!/usr/bin/env python3
"""Multi-modal memory support for WorkBuddy."""

from __future__ import annotations

import argparse
import base64
import json
import os
import tempfile
from pathlib import Path
from typing import Any, Dict, Optional, Union

# Try to import optional dependencies
try:
    import cv2
    import numpy as np
    from PIL import Image
    has_cv2 = True
except ImportError:
    has_cv2 = False

try:
    import pytesseract
    has_tesseract = True
except ImportError:
    has_tesseract = False

try:
    import whisper
    has_whisper = True
except ImportError:
    has_whisper = False

try:
    from transformers import CLIPProcessor, CLIPModel
    import torch
    has_clip = True
except ImportError:
    has_clip = False


class MultiModalProcessor:
    """Process multi-modal content for memory storage."""
    
    def __init__(self):
        self.clip_model = None
        self.clip_processor = None
        self.whisper_model = None
        self._load_models()
    
    def _load_models(self):
        """Load required models if available."""
        if has_clip:
            try:
                self.clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
                self.clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
            except Exception as e:
                print(f"Warning: Failed to load CLIP model: {e}")
        
        if has_whisper:
            try:
                self.whisper_model = whisper.load_model("base")
            except Exception as e:
                print(f"Warning: Failed to load Whisper model: {e}")
    
    def process_image(self, image_path: Union[str, Path]) -> Dict[str, Any]:
        """Process an image and extract content."""
        result = {
            "type": "image",
            "path": str(image_path),
            "content": {},
            "status": "success"
        }
        
        try:
            # Extract text using OCR
            if has_tesseract and has_cv2:
                img = cv2.imread(str(image_path))
                if img is not None:
                    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
                    text = pytesseract.image_to_string(gray)
                    if text.strip():
                        result["content"]["text"] = text.strip()
            
            # Extract image features using CLIP
            if has_clip:
                image = Image.open(image_path)
                inputs = self.clip_processor(images=image, return_tensors="pt")
                with torch.no_grad():
                    image_features = self.clip_model.get_image_features(**inputs)
                result["content"]["features"] = image_features.numpy().tolist()
        except Exception as e:
            result["status"] = "error"
            result["error"] = str(e)
        
        return result
    
    def process_audio(self, audio_path: Union[str, Path]) -> Dict[str, Any]:
        """Process an audio file and extract content."""
        result = {
            "type": "audio",
            "path": str(audio_path),
            "content": {},
            "status": "success"
        }
        
        try:
            if has_whisper:
                audio = whisper.load_audio(str(audio_path))
                audio = whisper.pad_or_trim(audio)
                mel = whisper.log_mel_spectrogram(audio).to(self.whisper_model.device)
                _, probs = self.whisper_model.detect_language(mel)
                options = whisper.DecodingOptions()
                text = whisper.decode(self.whisper_model, mel, options)
                result["content"]["text"] = text.text
                result["content"]["language"] = max(probs, key=probs.get)
        except Exception as e:
            result["status"] = "error"
            result["error"] = str(e)
        
        return result
    
    def process_video(self, video_path: Union[str, Path]) -> Dict[str, Any]:
        """Process a video file and extract content."""
        result = {
            "type": "video",
            "path": str(video_path),
            "content": {
                "frames": []
            },
            "status": "success"
        }
        
        try:
            if has_cv2:
                cap = cv2.VideoCapture(str(video_path))
                frame_count = 0
                while cap.isOpened() and frame_count < 10:  # Limit to 10 frames
                    ret, frame = cap.read()
                    if not ret:
                        break
                    
                    # Process every 30th frame
                    if frame_count % 30 == 0:
                        # Save frame to temporary file
                        with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
                            cv2.imwrite(tmp.name, frame)
                            frame_path = tmp.name
                        
                        # Process frame
                        frame_result = self.process_image(frame_path)
                        result["content"]["frames"].append(frame_result)
                        
                        # Clean up
                        os.unlink(frame_path)
                    
                    frame_count += 1
                
                cap.release()
        except Exception as e:
            result["status"] = "error"
            result["error"] = str(e)
        
        return result
    
    def process_pdf(self, pdf_path: Union[str, Path]) -> Dict[str, Any]:
        """Process a PDF file and extract content."""
        result = {
            "type": "pdf",
            "path": str(pdf_path),
            "content": {
                "text": ""
            },
            "status": "success"
        }
        
        try:
            # Try to import PyPDF2
            try:
                import PyPDF2
                has_pypdf2 = True
            except ImportError:
                has_pypdf2 = False
            
            if has_pypdf2:
                with open(pdf_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page_num in range(len(reader.pages)):
                        page = reader.pages[page_num]
                        text += page.extract_text() + "\n"
                    result["content"]["text"] = text.strip()
            else:
                result["status"] = "warning"
                result["warning"] = "PyPDF2 not available, cannot extract text from PDF"
        except Exception as e:
            result["status"] = "error"
            result["error"] = str(e)
        
        return result
    
    def process_document(self, doc_path: Union[str, Path]) -> Dict[str, Any]:
        """Process a document file and extract content."""
        result = {
            "type": "document",
            "path": str(doc_path),
            "content": {
                "text": ""
            },
            "status": "success"
        }
        
        try:
            extension = Path(doc_path).suffix.lower()
            
            if extension == ".txt":
                # Plain text file
                with open(doc_path, "r", encoding="utf-8", errors="ignore") as f:
                    result["content"]["text"] = f.read()
            elif extension in [".doc", ".docx"]:
                # Word document
                try:
                    import docx
                    has_docx = True
                except ImportError:
                    has_docx = False
                
                if has_docx:
                    doc = docx.Document(doc_path)
                    text = ""
                    for para in doc.paragraphs:
                        text += para.text + "\n"
                    result["content"]["text"] = text.strip()
                else:
                    result["status"] = "warning"
                    result["warning"] = "python-docx not available, cannot extract text from Word document"
            elif extension == ".rtf":
                # RTF file
                # Simple RTF to text conversion
                import re
                with open(doc_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
                    # Remove RTF control characters
                    text = re.sub(r"\\[a-z]{1,20}", "", text)
                    text = re.sub(r"\\'[0-9a-fA-F]{2}", "", text)
                    result["content"]["text"] = text.strip()
        except Exception as e:
            result["status"] = "error"
            result["error"] = str(e)
        
        return result
    
    def process_spreadsheet(self, spreadsheet_path: Union[str, Path]) -> Dict[str, Any]:
        """Process a spreadsheet file and extract content."""
        result = {
            "type": "spreadsheet",
            "path": str(spreadsheet_path),
            "content": {
                "sheets": []
            },
            "status": "success"
        }
        
        try:
            extension = Path(spreadsheet_path).suffix.lower()
            
            if extension == ".csv":
                # CSV file
                import csv
                with open(spreadsheet_path, "r", encoding="utf-8", errors="ignore") as f:
                    reader = csv.reader(f)
                    rows = list(reader)
                    result["content"]["sheets"].append({
                        "name": "Sheet1",
                        "rows": rows
                    })
            elif extension in [".xlsx", ".xls"]:
                # Excel file
                try:
                    import openpyxl
                    has_openpyxl = True
                except ImportError:
                    has_openpyxl = False
                
                if has_openpyxl:
                    wb = openpyxl.load_workbook(spreadsheet_path)
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        rows = []
                        for row in sheet.iter_rows(values_only=True):
                            rows.append(list(row))
                        result["content"]["sheets"].append({
                            "name": sheet_name,
                            "rows": rows
                        })
                else:
                    result["status"] = "warning"
                    result["warning"] = "openpyxl not available, cannot extract data from Excel file"
        except Exception as e:
            result["status"] = "error"
            result["error"] = str(e)
        
        return result
    
    def process_presentation(self, presentation_path: Union[str, Path]) -> Dict[str, Any]:
        """Process a presentation file and extract content."""
        result = {
            "type": "presentation",
            "path": str(presentation_path),
            "content": {
                "slides": []
            },
            "status": "success"
        }
        
        try:
            # Try to import python-pptx
            try:
                from pptx import Presentation
                has_pptx = True
            except ImportError:
                has_pptx = False
            
            if has_pptx:
                prs = Presentation(presentation_path)
                for i, slide in enumerate(prs.slides):
                    slide_content = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_content += shape.text + "\n"
                    result["content"]["slides"].append({
                        "index": i,
                        "content": slide_content.strip()
                    })
            else:
                result["status"] = "warning"
                result["warning"] = "python-pptx not available, cannot extract content from presentation"
        except Exception as e:
            result["status"] = "error"
            result["error"] = str(e)
        
        return result
    
    def process_file(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Process a file based on its extension."""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        if extension in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff"]:
            return self.process_image(file_path)
        elif extension in [".mp3", ".wav", ".ogg", ".flac", ".m4a", ".wma", ".aac"]:
            return self.process_audio(file_path)
        elif extension in [".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv", ".webm"]:
            return self.process_video(file_path)
        elif extension in [".pdf"]:
            return self.process_pdf(file_path)
        elif extension in [".doc", ".docx", ".txt", ".rtf"]:
            return self.process_document(file_path)
        elif extension in [".xlsx", ".xls", ".csv"]:
            return self.process_spreadsheet(file_path)
        elif extension in [".pptx", ".ppt"]:
            return self.process_presentation(file_path)
        else:
            return {
                "type": "unknown",
                "path": str(file_path),
                "status": "error",
                "error": "Unsupported file type"
            }
    
    def create_memory_from_multimodal(self, file_path: Union[str, Path], metadata: Optional[Dict[str, Any]] = None) -> str:
        """Create a memory from a multi-modal file."""
        file_path = Path(file_path)
        processed = self.process_file(file_path)
        
        if processed["status"] == "error":
            return f"Error processing file: {processed.get('error', 'Unknown error')}"
        
        # Create markdown content
        frontmatter = {
            "title": file_path.stem,
            "type": "multimodal",
            "created": metadata.get("created", "") if metadata else "",
            "tags": metadata.get("tags", "") if metadata else "",
            "file_type": processed["type"]
        }
        
        frontmatter_str = "---\n"
        for key, value in frontmatter.items():
            if value:
                frontmatter_str += f"{key}: {value}\n"
        frontmatter_str += "---\n\n"
        
        content = f"# {file_path.stem}\n\n"
        content += f"**File Path:** {processed['path']}\n"
        content += f"**File Type:** {processed['type']}\n\n"
        
        if processed.get("warning"):
            content += f"**Warning:** {processed['warning']}\n\n"
        
        if "text" in processed.get("content", {}):
            content += "## Extracted Text\n\n"
            # Limit text length to avoid overly large memories
            text = processed["content"]["text"]
            if len(text) > 1000:
                text = text[:1000] + "... (truncated)"
            content += text + "\n\n"
        
        if processed["type"] == "audio" and "language" in processed.get("content", {}):
            content += f"**Detected Language:** {processed['content']['language']}\n\n"
        
        if processed["type"] == "video" and "frames" in processed.get("content", {}):
            content += f"**Processed Frames:** {len(processed['content']['frames'])}\n\n"
        
        if processed["type"] == "spreadsheet" and "sheets" in processed.get("content", {}):
            content += "## Spreadsheet Data\n\n"
            for sheet in processed["content"]["sheets"]:
                content += f"### {sheet['name']}\n\n"
                # Convert rows to markdown table
                if sheet['rows']:
                    # Header row
                    header = sheet['rows'][0]
                    content += "| " + " | ".join(str(cell) for cell in header) + " |\n"
                    content += "| " + " | ".join(["---"] * len(header)) + " |\n"
                    # Data rows (limit to first 10 rows)
                    for row in sheet['rows'][1:11]:
                        content += "| " + " | ".join(str(cell) for cell in row) + " |\n"
                    if len(sheet['rows']) > 10:
                        content += f"... and {len(sheet['rows']) - 10} more rows\n\n"
                else:
                    content += "Empty sheet\n\n"
        
        if processed["type"] == "presentation" and "slides" in processed.get("content", {}):
            content += "## Presentation Slides\n\n"
            for slide in processed["content"]["slides"]:
                content += f"### Slide {slide['index'] + 1}\n\n"
                if slide['content']:
                    content += slide['content'] + "\n\n"
                else:
                    content += "No content\n\n"
        
        # Add base64 encoded thumbnail for images
        if processed["type"] == "image" and has_cv2:
            try:
                img = cv2.imread(str(file_path))
                if img is not None:
                    # Resize image for thumbnail
                    height, width = img.shape[:2]
                    max_dim = 300
                    if max(height, width) > max_dim:
                        scale = max_dim / max(height, width)
                        img = cv2.resize(img, (int(width * scale), int(height * scale)))
                    
                    # Convert to base64
                    _, buffer = cv2.imencode('.jpg', img)
                    img_base64 = base64.b64encode(buffer).decode('utf-8')
                    content += "## Thumbnail\n\n"
                    content += f"![Thumbnail](data:image/jpeg;base64,{img_base64})\n"
            except Exception as e:
                print(f"Error creating thumbnail: {e}")
        
        return frontmatter_str + content


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--file", required=True, help="Path to the multi-modal file")
    parser.add_argument("--output", help="Path to save the memory file")
    parser.add_argument("--json", action="store_true", help="Output JSON instead of markdown")
    args = parser.parse_args()
    
    processor = MultiModalProcessor()
    
    if args.json:
        result = processor.process_file(args.file)
        print(json.dumps(result, indent=2, ensure_ascii=False))
    else:
        memory_content = processor.create_memory_from_multimodal(args.file)
        
        if args.output:
            output_path = Path(args.output)
            output_path.write_text(memory_content, encoding="utf-8")
            print(f"Memory saved to: {output_path}")
        else:
            print(memory_content)
    
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
